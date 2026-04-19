import os
import chromadb
from pypdf import PdfReader
from pptx import Presentation

# Connect to Database
client = chromadb.PersistentClient('./RAG')
collection = client.get_or_create_collection('college_notes')

def process_pdf(file_path, course_code, filename):
    documents = []
    metadatas = []
    ids = []
    

    doc=PdfReader(file_path)
    for i,page in enumerate(doc.pages):
        text=page.extract_text()

        if text!="":
            documents.append(text)
            metadatas.append({"course":course_code,"source_file":filename,"file_type":"pdf","location":f"Page {i+1}"})
            ids.append(f"{filename}_page_{i+1}")


    return documents, metadatas, ids

def process_pptx(file_path, course_code, filename):
    documents = []
    metadatas = []
    ids = []
    
    prs = Presentation(file_path)
    
    # Loop through the slides
    for i, slide in enumerate(prs.slides):
        slide_text = []
        
        # Loop through every object (shape) on the slide
        for shape in slide.shapes:
            # Check if the object actually has text inside it
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        # Mash all the text boxes together into one giant string for the slide
        text = " ".join(slide_text)

        if text!="":
            documents.append(text)
            metadatas.append({"course":course_code,"source_file":filename,"file_type":".pptx","location":f"Slide {i+1}"})
            ids.append(f"{filename}_slide_no_{i+1}")
        
    return documents, metadatas, ids

# --- Main Ingestion Router ---
NOTES_FOLDER = "./notes"
COURSE_CODE = "OS2202"

all_docs, all_metadatas, all_ids = [], [], []

if not os.path.exists(NOTES_FOLDER):
    os.makedirs(NOTES_FOLDER)
    print(f"Created '{NOTES_FOLDER}' directory. Please place your course folders inside it.")

for course_folder in os.listdir(NOTES_FOLDER):
    course_path = os.path.join(NOTES_FOLDER,course_folder)

    if not os.path.isdir(course_path):
        continue

    COURSE_CODE = course_folder
    print(f"\n Scanning course folder: {COURSE_CODE}")

    for filename in os.listdir(course_path):
        if filename.startswith("~$"):
            continue
        file_path = os.path.join(course_path,filename)
        print(f"I see a file: {filename}")

        if filename.endswith(".pdf"):
            print(f"Extracting PDF: {filename}")
            docs, metas, doc_ids = process_pdf(file_path,COURSE_CODE,filename)
            all_docs.extend(docs)
            all_metadatas.extend(metas)
            all_ids.extend(doc_ids)
        
        elif filename.endswith(".pptx"):
            docs,metas,doc_ids = process_pptx(file_path,COURSE_CODE,filename)
            all_docs.extend(docs)
            all_metadatas.extend(metas)
            all_ids.extend(doc_ids)
        
        
        
# Finally, upload to ChromaDB
if all_docs:
    collection.add(documents=all_docs, metadatas=all_metadatas, ids=all_ids)
    print(f"Ingestion Complete! Uploaded {len(all_docs)} chunks.")
else:
    print("FAILED: No text was found to upload! Check your folder.")

